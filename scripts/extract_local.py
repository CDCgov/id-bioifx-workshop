#!/usr/bin/env python3
"""Extract text from local PDFs and PPTX files in the workshop repository.

Walks presentations/ and assets/pdfs/ for PDF and PPTX files, extracts
readable text, and writes it to _search/cache/<slug>.txt.

Usage:
    python scripts/extract_local.py          # incremental (skips if cache is newer)
    python scripts/extract_local.py --force  # re-extract everything
"""

import argparse
import sys
from pathlib import Path

# ---------------------------------------------------------------------------
# Text extraction helpers
# ---------------------------------------------------------------------------


def _extract_pdf(path: Path) -> str:
    from pdfminer.high_level import extract_text
    return extract_text(str(path))


def _extract_pptx(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(path))
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)
    return "\n".join(texts)


EXTRACTORS = {
    ".pdf": _extract_pdf,
    ".pptx": _extract_pptx,
}


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--force", action="store_true",
                        help="Re-extract even if cache is newer than source")
    args = parser.parse_args()

    repo_root = Path(__file__).resolve().parent.parent
    cache_dir = repo_root / "_search" / "cache"
    cache_dir.mkdir(parents=True, exist_ok=True)

    # Directories to scan for extractable files
    scan_dirs = [
        repo_root / "presentations",
        repo_root / "assets" / "pdfs",
    ]

    extracted = 0
    skipped = 0
    errors = 0

    for scan_dir in scan_dirs:
        if not scan_dir.is_dir():
            continue
        for filepath in sorted(scan_dir.rglob("*")):
            suffix = filepath.suffix.lower()
            if suffix not in EXTRACTORS:
                continue

            slug = filepath.stem
            cache_file = cache_dir / f"{slug}.txt"

            # Skip if cache is newer than source (incremental mode)
            if not args.force and cache_file.exists():
                if cache_file.stat().st_mtime >= filepath.stat().st_mtime:
                    skipped += 1
                    continue

            extractor = EXTRACTORS[suffix]
            try:
                text = extractor(filepath)
                cache_file.write_text(text, encoding="utf-8")
                extracted += 1
                print(f"  extracted: {filepath.relative_to(repo_root)} ({len(text)} chars)")
            except Exception as exc:
                print(f"  ERROR extracting {filepath}: {exc}", file=sys.stderr)
                errors += 1

    print(f"\nDone: {extracted} extracted, {skipped} skipped (cached), {errors} errors")
    return 1 if errors > 0 and extracted == 0 else 0


if __name__ == "__main__":
    sys.exit(main())
